"""Generate a Microsoft-branded PowerPoint deck from a YAML slide plan.

Usage:
    # Standalone — auto-discovers slides.yaml in same directory:
    python generate_pptx.py

    # Explicit paths:
    python generate_pptx.py --plan-file slides.yaml --output-file deck.pptx

    # With a template:
    python generate_pptx.py --plan-file slides.yaml --output-file deck.pptx --template template.potx
"""

import argparse
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Microsoft brand defaults
_BLUE = RGBColor(0x00, 0x78, 0xD4)
_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
_ACCENT_GRAY = RGBColor(0x60, 0x60, 0x60)

_FONT_FAMILY = "Segoe UI"
_TITLE_SIZE = Pt(28)
_SUBTITLE_SIZE = Pt(18)
_BODY_SIZE = Pt(18)
_SMALL_SIZE = Pt(14)
_TABLE_HEADER_SIZE = Pt(14)
_TABLE_BODY_SIZE = Pt(13)

_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)
_DEFAULT_OUTPUT_NAME = "presentation.pptx"
_LAYOUT_TITLE_ONLY = "Title Only"
_LAYOUT_TWO_COL_SUBHEADS = "2-column_Text_with Subheads"
_ASK_HEADING = "The Ask"


def _set_font(run, size, color=_DARK_TEXT, bold=False, font_family=_FONT_FAMILY):
    """Apply font properties to a text run."""
    run.font.name = font_family
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold


def _add_text_box(slide, left, top, width, height, text, size, color=_DARK_TEXT,
                  bold=False, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a text box with a single styled text run."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.paragraphs[0].alignment = alignment
    run = tf.paragraphs[0].add_run()
    run.text = text
    _set_font(run, size, color, bold)
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    return txbox


def _add_bullets(slide, left, top, width, height, bullets, size=_BODY_SIZE,
                 color=_DARK_TEXT):
    """Add a text box with bullet-pointed paragraphs."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, bullet in enumerate(bullets):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.space_after = Pt(8)
        para.level = 0
        run = para.add_run()
        run.text = bullet
        _set_font(run, size, color)

    return txbox


def _add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    if tf is not None:
        tf.text = notes_text


# ---------------------------------------------------------------------------
# Template configuration and layout helpers
# ---------------------------------------------------------------------------

# Default layout name preferences (used when no template.yaml exists)
_DEFAULT_LAYOUT_NAMES = {
    "title": ["Title Slide"],
    "section": ["Section Header"],
    "content": ["Title and Content"],
    "two-column": ["Two Content", "Comparison"],
    "two_column": ["Two Content", "Comparison"],
    "table": [_LAYOUT_TITLE_ONLY],
    "summary": [_LAYOUT_TITLE_ONLY],
    "ask": [_LAYOUT_TWO_COL_SUBHEADS, _LAYOUT_TITLE_ONLY],
    "blank": ["Blank"],
}

# Default placeholder mappings (standard PowerPoint placeholder indices)
_DEFAULT_PLACEHOLDERS = {
    "title": {"title": 0, "subtitle": 1},
    "section": {"title": 0, "subtitle": 1},
    "content": {"title": 0, "body": 1},
    "two-column": {"title": 0, "left_body": 1, "right_body": 2},
    "two_column": {"title": 0, "left_body": 1, "right_body": 2},
    "table": {"title": 0},
    "summary": {"title": 0},
    "ask": {"title": 0},
    "blank": {"title": 0},
}


def _load_template_config(template_path: str) -> dict | None:
    """Load a template.yaml sidecar config file if it exists.

    Looks for a .yaml file with the same base name as the template,
    next to the template file. Returns the parsed config or None.
    """
    if not template_path:
        return None

    config_path = Path(template_path).with_suffix(".yaml")
    if not config_path.exists():
        # Also try .yml
        config_path = Path(template_path).with_suffix(".yml")
    if not config_path.exists():
        return None

    with open(config_path, "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)

    print(f"Loaded template config: {config_path}")
    return config


def _get_layout_name(slide_type: str, template_config: dict | None) -> list[str]:
    """Get layout name(s) to search for the given slide type."""
    if template_config and "layouts" in template_config:
        name = template_config["layouts"].get(slide_type)
        if name:
            return [name] if isinstance(name, str) else name
    return _DEFAULT_LAYOUT_NAMES.get(slide_type, ["Blank"])


def _get_placeholder_idx(slide_type: str, role: str,
                         template_config: dict | None) -> int | None:
    """Get the placeholder index for a content role on a slide type."""
    if template_config and "placeholders" in template_config:
        type_map = template_config["placeholders"].get(slide_type, {})
        idx = type_map.get(role)
        if idx is not None:
            return idx
    return _DEFAULT_PLACEHOLDERS.get(slide_type, {}).get(role)


def _find_layout(prs, slide_type, has_template, template_config=None):
    """Find the best slide layout for the given type.

    When a template is loaded, tries named layouts from template config
    using prefix matching, then falls back to the first available layout.
    Without a template, uses the Blank layout for manual positioning.
    """
    if not has_template:
        idx = min(6, len(prs.slide_layouts) - 1)
        return prs.slide_layouts[idx]

    # Try name-based matching (case-insensitive, prefix match)
    preferred = _get_layout_name(slide_type, template_config)
    for name in preferred:
        for layout in prs.slide_layouts:
            layout_lower = layout.name.lower().strip()
            name_lower = name.lower().strip()
            if layout_lower == name_lower or layout_lower.startswith(name_lower):
                return layout

    # Fallback to first layout with a title placeholder
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == 0:
                return layout

    return prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]


def _try_set_placeholder(slide, idx, text):
    """Set text on a placeholder by index. Returns True if successful."""
    try:
        ph = slide.placeholders[idx]
        ph.text = text
        return True
    except (KeyError, IndexError):
        return False


def _try_set_body_bullets(slide, idx, bullets):
    """Fill a body placeholder with bullet points. Returns True if successful."""
    try:
        ph = slide.placeholders[idx]
        tf = ph.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.space_after = Pt(8)
            run = para.add_run()
            run.text = bullet
        return True
    except (KeyError, IndexError):
        return False


def _cleanup_empty_placeholders(slide):
    """Remove placeholder shapes that have no content to avoid empty boxes."""
    to_remove = []
    for ph in slide.placeholders:
        if hasattr(ph, "text") and not ph.text.strip():
            to_remove.append(ph)
    for ph in to_remove:
        elem = ph._element
        elem.getparent().remove(elem)


def _build_title_slide(prs, slide_data, has_template=False, template_config=None):
    """Build a title slide."""
    layout = _find_layout(prs, "title", has_template, template_config)
    slide = prs.slides.add_slide(layout)

    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")

    if has_template:
        title_idx = _get_placeholder_idx("title", "title", template_config)
        sub_idx = _get_placeholder_idx("title", "subtitle", template_config)
        if title_idx is not None:
            _try_set_placeholder(slide, title_idx, title)
        if sub_idx is not None:
            _try_set_placeholder(slide, sub_idx, subtitle)
    else:
        # Manual layout with accent bar
        from pptx.enum.shapes import MSO_SHAPE
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     _SLIDE_WIDTH, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _BLUE
        bar.line.fill.background()

        _add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
                      title, Pt(40), _DARK_TEXT, bold=True, alignment=PP_ALIGN.LEFT)
        if subtitle:
            _add_text_box(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.8),
                          subtitle, _SUBTITLE_SIZE, _ACCENT_GRAY)

    # Presenter and date — always manual (no standard placeholder for this)
    meta_parts = []
    if slide_data.get("presenter"):
        meta_parts.append(slide_data["presenter"])
    if slide_data.get("date"):
        meta_parts.append(slide_data["date"])
    if meta_parts:
        _add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
                      " | ".join(meta_parts), _SMALL_SIZE, _ACCENT_GRAY)

    if has_template:
        _cleanup_empty_placeholders(slide)
    _add_speaker_notes(slide, slide_data.get("speaker_notes"))
    return slide


def _build_section_slide(prs, slide_data, has_template=False, template_config=None):
    """Build a section divider slide."""
    layout = _find_layout(prs, "section", has_template, template_config)
    slide = prs.slides.add_slide(layout)

    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")

    if has_template:
        title_idx = _get_placeholder_idx("section", "title", template_config)
        sub_idx = _get_placeholder_idx("section", "subtitle", template_config)
        if title_idx is not None:
            _try_set_placeholder(slide, title_idx, title)
        if subtitle and sub_idx is not None:
            _try_set_placeholder(slide, sub_idx, subtitle)
    else:
        # Manual blue background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _BLUE

        _add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                      title, Pt(36), _WHITE, bold=True,
                      alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        if subtitle:
            _add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
                          subtitle, _SUBTITLE_SIZE, _WHITE)

    if has_template:
        _cleanup_empty_placeholders(slide)
    _add_speaker_notes(slide, slide_data.get("speaker_notes"))
    return slide


def _build_content_slide(prs, slide_data, has_template=False, template_config=None):
    """Build a title + bullets slide."""
    layout = _find_layout(prs, "content", has_template, template_config)
    slide = prs.slides.add_slide(layout)

    title = slide_data.get("title", "")
    bullets = slide_data.get("bullets", [])

    if has_template:
        title_idx = _get_placeholder_idx("content", "title", template_config)
        body_idx = _get_placeholder_idx("content", "body", template_config)
        if title_idx is not None:
            _try_set_placeholder(slide, title_idx, title)
        if bullets:
            placed = False
            if body_idx is not None:
                placed = _try_set_body_bullets(slide, body_idx, bullets)
            if not placed:
                _add_bullets(slide, Inches(0.8), Inches(1.7),
                             Inches(11), Inches(5.0), bullets)
    else:
        _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.9),
                      title, _TITLE_SIZE, _DARK_TEXT, bold=True)
        from pptx.enum.shapes import MSO_SHAPE
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35),
                                      Inches(1.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = _BLUE
        line.line.fill.background()
        if bullets:
            _add_bullets(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(5.0),
                         bullets)

    if has_template:
        _cleanup_empty_placeholders(slide)
    _add_speaker_notes(slide, slide_data.get("speaker_notes"))
    return slide


def _place_two_column_content(slide, slide_type, left_data, right_data,
                              template_config):
    """Place left/right column content into template placeholders."""
    lh_idx = _get_placeholder_idx(slide_type, "left_heading", template_config)
    lb_idx = _get_placeholder_idx(slide_type, "left_body", template_config)
    rh_idx = _get_placeholder_idx(slide_type, "right_heading", template_config)
    rb_idx = _get_placeholder_idx(slide_type, "right_body", template_config)

    left_heading = left_data.get("heading", "")
    left_body = "\n".join(left_data.get("bullets", []))
    right_heading = right_data.get("heading", "")
    right_body = "\n".join(right_data.get("bullets", []))

    placed = True
    for idx, text in [(lh_idx, left_heading), (lb_idx, left_body),
                      (rh_idx, right_heading), (rb_idx, right_body)]:
        if idx is not None:
            placed = _try_set_placeholder(slide, idx, text) and placed
    return placed


def _build_two_column_slide(prs, slide_data, has_template=False, template_config=None):
    """Build a two-column comparison slide."""
    layout = _find_layout(prs, "two-column", has_template, template_config)
    slide = prs.slides.add_slide(layout)

    title = slide_data.get("title", "")
    left_data = slide_data.get("left", {})
    right_data = slide_data.get("right", {})

    if has_template:
        title_idx = _get_placeholder_idx("two-column", "title", template_config)
        if title_idx is not None:
            _try_set_placeholder(slide, title_idx, title)
        if not _place_two_column_content(slide, "two-column", left_data,
                                         right_data, template_config):
            _build_manual_column(slide, left_data, Inches(0.8))
            _build_manual_column(slide, right_data, Inches(6.8))
    else:
        _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.9),
                      title, _TITLE_SIZE, _DARK_TEXT, bold=True)
        from pptx.enum.shapes import MSO_SHAPE
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35),
                                      Inches(1.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = _BLUE
        line.line.fill.background()
        _build_manual_column(slide, left_data, Inches(0.8))
        _build_manual_column(slide, right_data, Inches(6.8))

    if has_template:
        _cleanup_empty_placeholders(slide)
    _add_speaker_notes(slide, slide_data.get("speaker_notes"))
    return slide


def _column_to_text(col_data):
    """Convert a column dict (heading + bullets) to a single text string."""
    parts = []
    heading = col_data.get("heading", "")
    if heading:
        parts.append(heading)
    for bullet in col_data.get("bullets", []):
        parts.append(f"  {bullet}")
    return "\n".join(parts)


def _build_manual_column(slide, col_data, left_offset):
    """Build a column manually with heading + bullets."""
    col_width = Inches(5.2)
    heading = col_data.get("heading", "")
    if heading:
        _add_text_box(slide, left_offset, Inches(1.7), col_width, Inches(0.6),
                      heading, Pt(20), _BLUE, bold=True)
    bullets = col_data.get("bullets", [])
    if bullets:
        _add_bullets(slide, left_offset, Inches(2.4), col_width, Inches(4.3),
                     bullets)


def _style_cell_font(cell, size, color, bold=False):
    """Apply font styling to all runs in a table cell."""
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            _set_font(run, size, color, bold)


def _style_table_header(table, headers):
    """Style the header row of a table."""
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _BLUE
        _style_cell_font(cell, _TABLE_HEADER_SIZE, _WHITE, bold=True)


def _style_table_rows(table, rows):
    """Style the data rows of a table with alternating shading."""
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _LIGHT_GRAY if row_idx % 2 == 0 else _WHITE
            _style_cell_font(cell, _TABLE_BODY_SIZE, _DARK_TEXT)


def _build_table_slide(prs, slide_data, has_template=False, template_config=None):
    """Build a slide with a data table."""
    layout = _find_layout(prs, "table", has_template, template_config)
    slide = prs.slides.add_slide(layout)

    title = slide_data.get("title", "")
    if has_template:
        title_idx = _get_placeholder_idx("table", "title", template_config)
        if title_idx is not None:
            _try_set_placeholder(slide, title_idx, title)
    else:
        _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.9),
                      title, _TITLE_SIZE, _DARK_TEXT, bold=True)
        from pptx.enum.shapes import MSO_SHAPE
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35),
                                      Inches(1.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = _BLUE
        line.line.fill.background()

    table_data = slide_data.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers:
        _add_speaker_notes(slide, slide_data.get("speaker_notes"))
        return slide

    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    table_width = Inches(11.0)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(1.0), Inches(1.8),
        table_width, Inches(0.45 * num_rows)
    )
    table = table_shape.table

    _style_table_header(table, headers)
    _style_table_rows(table, rows)

    if has_template:
        _cleanup_empty_placeholders(slide)
    _add_speaker_notes(slide, slide_data.get("speaker_notes"))
    return slide


def _build_summary_slide(prs, slide_data, has_template=False, template_config=None):
    """Build a summary/takeaway slide with a headline metric."""
    layout = _find_layout(prs, "summary", has_template, template_config)
    slide = prs.slides.add_slide(layout)

    title = slide_data.get("title", "")
    headline = slide_data.get("headline", "")
    bullets = slide_data.get("bullets", [])

    if has_template:
        title_idx = _get_placeholder_idx("summary", "title", template_config)
        body_idx = _get_placeholder_idx("summary", "body", template_config)
        if title_idx is not None:
            _try_set_placeholder(slide, title_idx, headline or title)
        bullets_text = "\n".join(bullets)
        if body_idx is not None and not _try_set_placeholder(slide, body_idx, bullets_text) and bullets:
            _add_bullets(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(3.0),
                         bullets, _BODY_SIZE, _ACCENT_GRAY)
    else:
        _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.9),
                      title, _TITLE_SIZE, _DARK_TEXT, bold=True)
        if headline:
            _add_text_box(slide, Inches(1.0), Inches(2.0), Inches(11), Inches(1.2),
                          headline, Pt(36), _BLUE, bold=True,
                          alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if bullets:
            _add_bullets(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(3.0),
                         bullets, _BODY_SIZE, _ACCENT_GRAY)

    if has_template:
        _cleanup_empty_placeholders(slide)
    _add_speaker_notes(slide, slide_data.get("speaker_notes"))
    return slide


def _build_manual_ask(slide, slide_data):
    """Build ask slide content manually without template placeholders."""
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.9),
                  slide_data.get("title", _ASK_HEADING), _TITLE_SIZE, _DARK_TEXT, bold=True)
    from pptx.enum.shapes import MSO_SHAPE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35),
                                  Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = _BLUE
    line.line.fill.background()

    ask_text = slide_data.get("ask", "")
    if ask_text:
        _add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.0),
                      ask_text, Pt(22), _BLUE, bold=True)

    rationale = slide_data.get("rationale", [])
    if rationale:
        _add_text_box(slide, Inches(0.8), Inches(3.0), Inches(5), Inches(0.5),
                      "Why", Pt(16), _DARK_TEXT, bold=True)
        _add_bullets(slide, Inches(0.8), Inches(3.5), Inches(5), Inches(2.5),
                     rationale, Pt(16), _ACCENT_GRAY)

    next_steps = slide_data.get("next_steps", [])
    if next_steps:
        _add_text_box(slide, Inches(6.8), Inches(3.0), Inches(5), Inches(0.5),
                      "Next Steps", Pt(16), _DARK_TEXT, bold=True)
        _add_bullets(slide, Inches(6.8), Inches(3.5), Inches(5), Inches(2.5),
                     next_steps, Pt(16), _ACCENT_GRAY)


def _build_ask_slide(prs, slide_data, has_template=False, template_config=None):
    """Build a decision/ask slide."""
    layout = _find_layout(prs, "ask", has_template, template_config)
    slide = prs.slides.add_slide(layout)

    if has_template:
        title = slide_data.get("title", _ASK_HEADING)
        title_idx = _get_placeholder_idx("ask", "title", template_config)
        if title_idx is not None:
            _try_set_placeholder(slide, title_idx, title)

        ask_text = slide_data.get("ask", "")
        rationale = slide_data.get("rationale", [])
        next_steps = slide_data.get("next_steps", [])

        left_body = ask_text
        if rationale:
            left_body += "\n\nWhy:\n" + "\n".join(f"  {r}" for r in rationale)

        left_data = {"heading": _ASK_HEADING if ask_text else "",
                     "bullets": [left_body] if left_body else []}
        right_data = {"heading": "Next Steps" if next_steps else "",
                      "bullets": next_steps}

        _place_two_column_content(slide, "ask", left_data, right_data,
                                  template_config)
    else:
        _build_manual_ask(slide, slide_data)

    if has_template:
        _cleanup_empty_placeholders(slide)
    _add_speaker_notes(slide, slide_data.get("speaker_notes"))
    return slide


def _build_blank_slide(prs, slide_data, has_template=False, template_config=None):
    """Build a blank slide with only a title (for manual editing)."""
    layout = _find_layout(prs, "blank", has_template, template_config)
    slide = prs.slides.add_slide(layout)

    title = slide_data.get("title", "")
    if has_template:
        title_idx = _get_placeholder_idx("blank", "title", template_config)
        if title_idx is None or not _try_set_placeholder(slide, title_idx, title):
            _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.9),
                          title, _TITLE_SIZE, _DARK_TEXT, bold=True)
    else:
        _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.9),
                      title, _TITLE_SIZE, _DARK_TEXT, bold=True)

    if has_template:
        _cleanup_empty_placeholders(slide)
    _add_speaker_notes(slide, slide_data.get("speaker_notes"))
    return slide


_SLIDE_BUILDERS = {
    "title": _build_title_slide,
    "section": _build_section_slide,
    "content": _build_content_slide,
    "two-column": _build_two_column_slide,
    "two_column": _build_two_column_slide,
    "table": _build_table_slide,
    "summary": _build_summary_slide,
    "ask": _build_ask_slide,
    "blank": _build_blank_slide,
}


def _resolve_template(plan_path: str, plan: dict,
                      cli_template: str | None) -> str | None:
    """Resolve the template path from CLI arg or YAML plan field.

    Priority: CLI --template > YAML template field > None.
    The YAML template field can be an absolute path, a relative path
    (resolved against the plan file directory), or a short name that
    maps to a templates/ folder next to the plan file or next to this script.
    """
    if cli_template:
        resolved = Path(cli_template)
        if resolved.exists():
            return str(resolved)
        print(f"Warning: CLI template not found: {cli_template}")

    yaml_template = plan.get("template")
    if not yaml_template:
        return None

    # Try as-is (absolute or relative to cwd)
    candidate = Path(yaml_template)
    if candidate.exists():
        return str(candidate)

    plan_dir = Path(plan_path).resolve().parent
    script_dir = Path(__file__).resolve().parent

    # Search directories: plan file dir, script dir, skill templates dir
    search_dirs = [
        plan_dir,
        script_dir,
        plan_dir / "templates",
        script_dir.parent / "templates",
    ]

    for search_dir in search_dirs:
        # Try exact name, then with extensions appended
        for ext in ("", ".potx", ".pptx"):
            candidate = search_dir / f"{yaml_template}{ext}"
            if candidate.exists():
                return str(candidate)

    print(f"Note: Template '{yaml_template}' not found — using built-in defaults.")
    return None


def _load_template(template_path: str) -> Presentation:
    """Load a template file, converting .potx to .pptx if needed.

    python-pptx does not support .potx files directly. This function
    copies the .potx to a temp .pptx and patches the content type
    inside [Content_Types].xml so python-pptx can open it.
    """
    path = Path(template_path)

    if path.suffix.lower() != ".potx":
        return Presentation(str(path))

    # Copy to a temp .pptx and patch the content type
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp_path = tmp.name
    tmp.close()

    try:
        shutil.copy2(str(path), tmp_path)

        _POTX_CONTENT_TYPE = (
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.template.main+xml"
        )
        _PPTX_CONTENT_TYPE = (
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation.main+xml"
        )

        with zipfile.ZipFile(tmp_path, "r") as zin:
            content_types = zin.read("[Content_Types].xml").decode("utf-8")

        patched = content_types.replace(_POTX_CONTENT_TYPE, _PPTX_CONTENT_TYPE)

        # Rewrite [Content_Types].xml inside the zip
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            with zipfile.ZipFile(str(path), "r") as zin:
                for item in zin.infolist():
                    if item.filename == "[Content_Types].xml":
                        zout.writestr(item, patched)
                    else:
                        zout.writestr(item, zin.read(item.filename))

        return Presentation(tmp_path)
    finally:
        Path(tmp_path).unlink(missing_ok=True)


def _clear_sections(prs):
    """Remove pre-existing section list from a template.

    Templates often carry section groupings (e.g. "Title Slides",
    "Executive Slides") that cause all generated slides to pile up
    under the first section. Removing the sectionLst element lets
    PowerPoint treat generated slides as ungrouped.
    """
    from lxml import etree
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    presentation_elem = prs._element
    for section_lst in presentation_elem.findall(".//p:sectionLst", nsmap):
        section_lst.getparent().remove(section_lst)
    # Also check for the extended namespace used by newer templates
    nsmap_ext = {"p14": "http://schemas.microsoft.com/office/powerpoint/2010/main"}
    for section_lst in presentation_elem.findall(".//p14:sectionLst", nsmap_ext):
        section_lst.getparent().remove(section_lst)


def generate_presentation(plan_path: str, output_path: str,
                          template_path: str | None = None) -> str:
    """Generate a PowerPoint deck from a YAML slide plan.

    Args:
        plan_path: Path to YAML slide plan file.
        output_path: Output path for the .pptx file.
        template_path: Optional path to a .potx/.pptx template (CLI override).

    Returns:
        Status message.
    """
    plan_file = Path(plan_path)
    if not plan_file.exists():
        return f"Error: Plan file not found: {plan_path}"

    with open(plan_file, "r", encoding="utf-8") as f:
        plan = yaml.safe_load(f)

    if not plan or "slides" not in plan:
        return "Error: Plan file must contain a 'slides' list."

    # Resolve template: CLI arg > YAML field > None
    resolved_template = _resolve_template(plan_path, plan, template_path)

    # Create presentation
    if resolved_template:
        prs = _load_template(resolved_template)
        _clear_sections(prs)
        template_config = _load_template_config(resolved_template)
        print(f"Using template: {resolved_template}")
    else:
        prs = Presentation()
        template_config = None

    # Set slide dimensions
    aspect = plan.get("aspect_ratio", "16:9")
    if aspect == "4:3":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = _SLIDE_WIDTH
        prs.slide_height = _SLIDE_HEIGHT

    slides = plan.get("slides", [])
    _apply_metadata_defaults(plan, slides)

    has_template = bool(resolved_template)
    for slide_data in slides:
        _build_slide(prs, slide_data, has_template, template_config)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))

    return f"Generated {len(slides)}-slide presentation at {output_path}"


def _apply_metadata_defaults(plan, slides):
    """Propagate top-level presenter/date to title slides missing them."""
    top_presenter = plan.get("presenter", "")
    top_date = plan.get("date", "")
    for slide_data in slides:
        if slide_data.get("type", "content").lower().strip() == "title":
            slide_data.setdefault("presenter", top_presenter)
            slide_data.setdefault("date", top_date)


def _build_slide(prs, slide_data, has_template=False, template_config=None):
    """Dispatch a single slide to the appropriate builder."""
    slide_type = slide_data.get("type", "content").lower().strip()
    builder = _SLIDE_BUILDERS.get(slide_type)
    if builder is None:
        print(f"Warning: Unknown slide type '{slide_type}', using 'content'.")
        builder = _build_content_slide
    builder(prs, slide_data, has_template, template_config)


def _resolve_defaults():
    """Auto-discover slides.yaml and output path relative to this script."""
    script_dir = Path(__file__).resolve().parent
    candidates = ["slides.yaml", "slides.yml"]
    for name in candidates:
        candidate = script_dir / name
        if candidate.exists():
            return str(candidate), str(script_dir / _DEFAULT_OUTPUT_NAME)
    # Also check current working directory
    cwd = Path.cwd()
    for name in candidates:
        candidate = cwd / name
        if candidate.exists():
            return str(candidate), str(cwd / _DEFAULT_OUTPUT_NAME)
    return None, None


def main():
    parser = argparse.ArgumentParser(
        description="Generate a Microsoft-branded PowerPoint deck from a YAML slide plan."
    )
    parser.add_argument(
        "--plan-file", default=None,
        help="Path to the YAML slide plan file. "
             "If omitted, auto-discovers slides.yaml next to this script or in cwd."
    )
    parser.add_argument(
        "--output-file", default=None,
        help="Output path for the .pptx file. "
             "Defaults to presentation.pptx next to the plan file."
    )
    parser.add_argument(
        "--template", default=None,
        help="Optional path to a .potx or .pptx template file."
    )

    args = parser.parse_args()

    plan_file = args.plan_file
    output_file = args.output_file

    # Auto-discover if no plan file specified
    if not plan_file:
        plan_file, default_output = _resolve_defaults()
        if not plan_file:
            print(
                "Error: No slides.yaml found. Provide --plan-file "
                "or place slides.yaml next to this script.",
                file=sys.stderr,
            )
            sys.exit(1)
        if not output_file:
            output_file = default_output
        print(f"Auto-discovered plan: {plan_file}")

    # Default output next to plan file
    if not output_file:
        output_file = str(Path(plan_file).parent / _DEFAULT_OUTPUT_NAME)

    try:
        result = generate_presentation(plan_file, output_file, args.template)
        print(result)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
